#!/usr/bin/env python3
"""PPT 大纲 → PPTX 转换脚本（Designer Skill）"""
import sys
import argparse
import re
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


DELL_BLUE = RGBColor(0x00, 0x76, 0xCE)
DELL_DARK = RGBColor(0x12, 0x3D, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def outline_to_pptx(input_path: str, output_path: str, template: str = None):
    outline_file = Path(input_path)
    pptx_file = Path(output_path)
    pptx_file.parent.mkdir(parents=True, exist_ok=True)

    content = outline_file.read_text(encoding="utf-8")

    # Parse outline
    slides = []
    current = None

    for line in content.split('\n'):
        m = re.match(r'^##\s+(Slide\s+\d+:?\s*)?(.+)$', line)
        if m:
            if current:
                slides.append(current)
            current = {'title': m.group(2).strip(), 'points': []}
        elif current and line.strip().startswith('- '):
            current['points'].append(line.strip()[2:])
    if current:
        slides.append(current)

    # Build PPT
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # Title slide
    title_slide = prs.slides.add_slide(blank)
    title_shape = title_slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
    tf = title_shape.text_frame
    tf.text = "Dell Technologies"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = DELL_BLUE
    p.font.bold = True

    # Content slides
    for s in slides:
        slide = prs.slides.add_slide(blank)

        # Title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(1))
        tf = title_shape.text_frame
        tf.text = s['title']
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = DELL_BLUE
        p.font.bold = True

        # Content
        if s['points']:
            body_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12), Inches(5.5))
            tf = body_shape.text_frame
            tf.word_wrap = True
            for i, point in enumerate(s['points']):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(16)
                p.font.color.rgb = DELL_DARK
                p.space_after = Pt(8)

    prs.save(str(pptx_file))
    print(f"PPTX generated: {pptx_file} ({len(slides)+1} slides)")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Convert PPT outline to PPTX")
    parser.add_argument("--input", "-i", required=True, help="Input outline markdown")
    parser.add_argument("--output", "-o", required=True, help="Output PPTX path")
    parser.add_argument("--template", "-t", help="Optional template path")
    args = parser.parse_args()

    outline_to_pptx(args.input, args.output, args.template)